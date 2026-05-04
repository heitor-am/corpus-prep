"""Shared fixtures."""

from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def write_file(tmp_path: Path):
    """Helper for creating temp files with given content.

    Usage:
        path = write_file("test.txt", "content", encoding="utf-8")
    """

    def _write(name: str, content: str | bytes, encoding: str = "utf-8") -> Path:
        path = tmp_path / name
        if isinstance(content, bytes):
            path.write_bytes(content)
        else:
            path.write_text(content, encoding=encoding)
        return path

    return _write


@pytest.fixture
def make_native_pdf(tmp_path: Path):
    """Create a PDF with a real text layer using ReportLab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    def _make(name: str = "native.pdf", text: str = "Hello world.\nLine two.") -> Path:
        path = tmp_path / name
        c = canvas.Canvas(str(path), pagesize=A4)
        c.setFont("Helvetica", 12)
        y = 800
        for line in text.split("\n"):
            c.drawString(50, y, line)
            y -= 20
        c.showPage()
        c.save()
        return path

    return _make


@pytest.fixture
def make_blank_pdf(tmp_path: Path):
    """Create a PDF with no extractable text (just a geometric shape).

    Simulates a scanned PDF: PyMuPDF4LLM extracts zero text, triggering
    needs_ocr=true.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    def _make(name: str = "blank.pdf") -> Path:
        path = tmp_path / name
        c = canvas.Canvas(str(path), pagesize=A4)
        c.rect(50, 50, 500, 700, stroke=1, fill=0)
        c.showPage()
        c.save()
        return path

    return _make


@pytest.fixture
def make_docx(tmp_path: Path):
    """Create a small DOCX via python-docx."""

    def _make(name: str = "sample.docx", paragraphs: list[str] | None = None) -> Path:
        from docx import Document as DocxDocument

        doc = DocxDocument()
        for line in paragraphs or [
            "Title of the test document",
            "First paragraph with some content for parsing.",
            "Second paragraph adds more substance.",
        ]:
            doc.add_paragraph(line)
        path = tmp_path / name
        doc.save(str(path))
        return path

    return _make


@pytest.fixture
def make_pptx(tmp_path: Path):
    """Create a small PPTX via python-pptx."""

    def _make(name: str = "sample.pptx", slides: list[tuple[str, str]] | None = None) -> Path:
        from pptx import Presentation

        prs = Presentation()
        for title, body in slides or [
            ("Test deck", "Slide one body text"),
            ("Second slide", "Slide two body text"),
        ]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        path = tmp_path / name
        prs.save(str(path))
        return path

    return _make


@pytest.fixture
def make_image_only_pdf(tmp_path: Path):
    """Render text into an image and embed it in a PDF (no real text layer).

    Simulates a scanned document: PyMuPDF extracts ~zero chars, triggering
    needs_ocr=true; Docling fallback can OCR the image to recover the text.
    """

    def _make(name: str = "scanned.pdf", text: str = "Hello OCR world from corpus-prep") -> Path:
        from io import BytesIO

        from PIL import Image, ImageDraw
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas

        img = Image.new("RGB", (1200, 400), color="white")
        draw = ImageDraw.Draw(img)
        draw.text((40, 80), text, fill="black")
        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        path = tmp_path / name
        c = canvas.Canvas(str(path), pagesize=A4)
        c.drawImage(ImageReader(buf), 50, 400, width=500, height=200)
        c.showPage()
        c.save()
        return path

    return _make


@pytest.fixture
def make_text_image(tmp_path: Path):
    """Create a PNG with rendered text via PIL."""

    def _make(
        name: str = "sample.png",
        text: str = "Hello world from corpus-prep",
        size: tuple[int, int] = (600, 200),
    ) -> Path:
        from PIL import Image, ImageDraw

        img = Image.new("RGB", size, color="white")
        draw = ImageDraw.Draw(img)
        draw.text((20, 60), text, fill="black")
        path = tmp_path / name
        img.save(path)
        return path

    return _make
